"""Text-layer parsers for PPTX and PDF slide decks."""

from __future__ import annotations

import re
from collections.abc import Iterable
from io import BytesIO
from pathlib import Path
from typing import Any, BinaryIO

from .errors import (
    DocumentParseError,
    EncryptedDocumentError,
    ParserDependencyError,
    TextlessDocumentError,
    UnsupportedFileTypeError,
)
from .models import BlockKind, ParsedBlock, ParsedDeck, ParsedSlide, SourceFormat

DocumentSource = str | Path | bytes | bytearray | BinaryIO
_BULLET_RE = re.compile(r"^\s*(?:[-*•▪◦‣⁃]|\d+[.)]|[A-Za-z][.)])\s+")
_PAGE_NUMBER_RE = re.compile(r"^(?:page\s*)?\d+(?:\s*/\s*\d+)?$", re.IGNORECASE)


def parse_document(
    source: DocumentSource,
    *,
    filename: str | None = None,
    deck_title: str | None = None,
) -> ParsedDeck:
    """Route a supported file to its text parser using a trusted filename."""

    effective_name = filename
    if effective_name is None and isinstance(source, (str, Path)):
        effective_name = Path(source).name
    suffix = Path(effective_name or "").suffix.lower()
    if suffix == ".pptx":
        return parse_pptx(source, deck_title=deck_title)
    if suffix == ".pdf":
        return parse_pdf(source, deck_title=deck_title)
    if suffix == ".ppt":
        raise UnsupportedFileTypeError(
            "Legacy .ppt files are not supported; convert the deck to .pptx."
        )
    raise UnsupportedFileTypeError("Only text-layer .pptx and .pdf files are supported.")


def parse_pptx(
    source: DocumentSource,
    *,
    deck_title: str | None = None,
) -> ParsedDeck:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParserDependencyError("PPTX ingestion requires the 'python-pptx' package.") from exc

    try:
        presentation = Presentation(_library_source(source))
    except Exception as exc:
        raise DocumentParseError("The PPTX file could not be opened.") from exc

    inferred_title = deck_title or getattr(presentation.core_properties, "title", "")
    if not inferred_title and isinstance(source, (str, Path)):
        inferred_title = Path(source).stem

    parsed_slides: list[ParsedSlide] = []
    current_section = ""
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = getattr(slide.shapes, "title", None)
        title_shape_id = getattr(title_shape, "shape_id", None)
        title = _shape_text(title_shape) if title_shape is not None else ""
        blocks: list[ParsedBlock] = []

        shapes = sorted(
            _flatten_shapes(slide.shapes),
            key=lambda shape: (
                int(getattr(shape, "top", 0)),
                int(getattr(shape, "left", 0)),
                int(getattr(shape, "shape_id", 0)),
            ),
        )
        for shape in shapes:
            if title_shape_id is not None and getattr(shape, "shape_id", None) == title_shape_id:
                continue
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [_clean_inline(cell.text) for cell in row.cells]
                    row_text = " | ".join(cell for cell in cells if cell)
                    if row_text:
                        blocks.append(
                            ParsedBlock(
                                kind=BlockKind.TABLE_ROW_GROUP,
                                text=row_text,
                                reading_order=len(blocks),
                            )
                        )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            paragraphs = [
                paragraph
                for paragraph in shape.text_frame.paragraphs
                if _clean_inline(paragraph.text)
            ]
            if not paragraphs:
                continue
            is_bullet_group = (
                len(paragraphs) > 1
                or _shape_is_body_placeholder(shape)
                or any(_paragraph_is_bullet(paragraph) for paragraph in paragraphs)
            )
            if is_bullet_group:
                lines = []
                for paragraph in paragraphs:
                    indent = "  " * max(int(getattr(paragraph, "level", 0)), 0)
                    lines.append(f"{indent}- {_clean_inline(paragraph.text)}")
                text = "\n".join(lines)
                kind = BlockKind.BULLET_GROUP
            else:
                text = _clean_inline(paragraphs[0].text)
                kind = BlockKind.PARAGRAPH
            blocks.append(ParsedBlock(kind=kind, text=text, reading_order=len(blocks)))

        layout_name = str(getattr(getattr(slide, "slide_layout", None), "name", ""))
        if title and "section" in layout_name.casefold():
            current_section = title
        parsed_slides.append(
            ParsedSlide(
                number=slide_number,
                title=title,
                section=current_section,
                blocks=tuple(blocks),
            )
        )

    deck = ParsedDeck(
        title=(inferred_title or "Untitled deck").strip(),
        source_format=SourceFormat.PPTX,
        slides=tuple(parsed_slides),
    )
    _ensure_extractable_text(deck)
    return deck


def parse_pdf(
    source: DocumentSource,
    *,
    deck_title: str | None = None,
) -> ParsedDeck:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ParserDependencyError("PDF ingestion requires the 'pypdf' package.") from exc

    try:
        reader = PdfReader(_library_source(source))
        if reader.is_encrypted:
            try:
                decrypt_result = reader.decrypt("")
            except Exception as exc:
                raise EncryptedDocumentError("Encrypted PDFs are not supported.") from exc
            if not decrypt_result:
                raise EncryptedDocumentError("Encrypted PDFs are not supported.")
    except EncryptedDocumentError:
        raise
    except Exception as exc:
        raise DocumentParseError("The PDF file could not be opened.") from exc

    inferred_title = deck_title
    if not inferred_title and isinstance(source, (str, Path)):
        inferred_title = Path(source).stem
    if not inferred_title:
        metadata = getattr(reader, "metadata", None)
        inferred_title = getattr(metadata, "title", "") if metadata else ""

    parsed_slides: list[ParsedSlide] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            raise DocumentParseError(
                f"Could not extract text from PDF page {page_number}."
            ) from exc
        title, blocks = _parse_pdf_page_text(page_text)
        parsed_slides.append(
            ParsedSlide(
                number=page_number,
                title=title,
                blocks=blocks,
            )
        )

    deck = ParsedDeck(
        title=(inferred_title or "Untitled deck").strip(),
        source_format=SourceFormat.PDF,
        slides=tuple(parsed_slides),
    )
    _ensure_extractable_text(deck)
    return deck


def _library_source(source: DocumentSource) -> Any:
    if isinstance(source, (bytes, bytearray)):
        return BytesIO(bytes(source))
    return source


def _flatten_shapes(shapes: Iterable[Any]) -> Iterable[Any]:
    for shape in shapes:
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _flatten_shapes(nested)
        else:
            yield shape


def _shape_text(shape: Any) -> str:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(
        _clean_inline(paragraph.text)
        for paragraph in shape.text_frame.paragraphs
        if _clean_inline(paragraph.text)
    )


def _paragraph_is_bullet(paragraph: Any) -> bool:
    if int(getattr(paragraph, "level", 0)) > 0:
        return True
    paragraph_xml = getattr(paragraph, "_p", None)
    properties = getattr(paragraph_xml, "pPr", None)
    if properties is None:
        return False
    bullet_tags = {"buChar", "buAutoNum", "buBlip"}
    no_bullet = False
    for child in properties:
        local_name = str(child.tag).rsplit("}", 1)[-1]
        if local_name == "buNone":
            no_bullet = True
        if local_name in bullet_tags:
            return True
    return False if no_bullet else False


def _shape_is_body_placeholder(shape: Any) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        placeholder_type = shape.placeholder_format.type
    except (AttributeError, ValueError):
        return False
    return getattr(placeholder_type, "name", "") in {"BODY", "OBJECT"}


def _parse_pdf_page_text(text: str) -> tuple[str, tuple[ParsedBlock, ...]]:
    raw_lines = text.replace("\r\n", "\n").replace("\r", "\n").split("\n")
    lines = [_clean_inline(line) for line in raw_lines]
    while lines and not lines[0]:
        lines.pop(0)
    while lines and not lines[-1]:
        lines.pop()
    while lines and _PAGE_NUMBER_RE.fullmatch(lines[0]):
        lines.pop(0)
        while lines and not lines[0]:
            lines.pop(0)
    while lines and _PAGE_NUMBER_RE.fullmatch(lines[-1]):
        lines.pop()
        while lines and not lines[-1]:
            lines.pop()
    if not lines:
        return "", ()

    title = ""
    first_nonempty_index = next((i for i, line in enumerate(lines) if line), None)
    if first_nonempty_index is not None:
        candidate = lines[first_nonempty_index]
        if len(candidate) <= 240 and not _PAGE_NUMBER_RE.fullmatch(candidate):
            title = candidate
            lines = lines[first_nonempty_index + 1 :]

    blocks: list[ParsedBlock] = []
    paragraph_lines: list[str] = []
    bullet_lines: list[str] = []

    def flush_paragraph() -> None:
        if paragraph_lines:
            blocks.append(
                ParsedBlock(
                    kind=BlockKind.PARAGRAPH,
                    text=" ".join(paragraph_lines),
                    reading_order=len(blocks),
                )
            )
            paragraph_lines.clear()

    def flush_bullets() -> None:
        if bullet_lines:
            blocks.append(
                ParsedBlock(
                    kind=BlockKind.BULLET_GROUP,
                    text="\n".join(bullet_lines),
                    reading_order=len(blocks),
                )
            )
            bullet_lines.clear()

    for line in lines:
        if not line:
            flush_paragraph()
            flush_bullets()
        elif _BULLET_RE.match(line):
            flush_paragraph()
            bullet_lines.append(f"- {_BULLET_RE.sub('', line, count=1)}")
        else:
            flush_bullets()
            paragraph_lines.append(line)
    flush_paragraph()
    flush_bullets()
    return title, tuple(blocks)


def _clean_inline(text: str) -> str:
    return re.sub(r"\s+", " ", text or "").strip()


def _ensure_extractable_text(deck: ParsedDeck) -> None:
    if not any(
        slide.title.strip() or any(block.text.strip() for block in slide.blocks)
        for slide in deck.slides
    ):
        raise TextlessDocumentError(
            f"{deck.source_format.value.upper()} contains no extractable text; "
            "OCR/vision ingestion is not enabled."
        )
