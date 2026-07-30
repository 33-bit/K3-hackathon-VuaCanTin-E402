from io import BytesIO

import pytest

from app.ingestion import (
    BlockKind,
    EncryptedDocumentError,
    SourceFormat,
    TextlessDocumentError,
    UnsupportedFileTypeError,
    parse_document,
    parse_pdf,
    parse_pptx,
)


class _FakePage:
    def __init__(self, text: str | None) -> None:
        self.text = text

    def extract_text(self) -> str | None:
        return self.text


class _FakePdfReader:
    def __init__(
        self,
        pages: list[_FakePage],
        *,
        encrypted: bool = False,
        decrypt_result: int = 1,
    ) -> None:
        self.pages = pages
        self.is_encrypted = encrypted
        self.decrypt_result = decrypt_result
        self.metadata = None

    def decrypt(self, password: str) -> int:
        assert password == ""
        return self.decrypt_result


def test_parse_document_rejects_legacy_ppt_and_unknown_extensions() -> None:
    with pytest.raises(UnsupportedFileTypeError):
        parse_document(b"legacy", filename="slides.ppt")
    with pytest.raises(UnsupportedFileTypeError):
        parse_document(b"text", filename="slides.txt")


def test_pdf_parser_preserves_page_and_bullet_structure(monkeypatch) -> None:
    import pypdf

    fake_reader = _FakePdfReader(
        [
            _FakePage(
                "Giới thiệu Qdrant\nĐây là phần mở đầu.\n\n• Dense retrieval\n- Sparse retrieval\n"
            ),
            _FakePage("Trang hai\nNội dung tiếp theo"),
        ]
    )
    monkeypatch.setattr(pypdf, "PdfReader", lambda source: fake_reader)

    deck = parse_pdf(b"%PDF-fake", deck_title="RAG")

    assert deck.source_format is SourceFormat.PDF
    assert deck.title == "RAG"
    assert len(deck.slides) == 2
    assert deck.slides[0].title == "Giới thiệu Qdrant"
    assert [block.kind for block in deck.slides[0].blocks] == [
        BlockKind.PARAGRAPH,
        BlockKind.BULLET_GROUP,
    ]
    assert deck.slides[0].blocks[1].text == ("- Dense retrieval\n- Sparse retrieval")


def test_pdf_parser_raises_typed_error_for_textless_pdf(monkeypatch) -> None:
    import pypdf

    monkeypatch.setattr(
        pypdf,
        "PdfReader",
        lambda source: _FakePdfReader([_FakePage(None), _FakePage(" \n "), _FakePage("Page 3")]),
    )

    with pytest.raises(TextlessDocumentError) as exc_info:
        parse_pdf(b"%PDF-fake")

    assert exc_info.value.code == "unsupported_textless_document"


def test_pdf_parser_rejects_encrypted_pdf(monkeypatch) -> None:
    import pypdf

    monkeypatch.setattr(
        pypdf,
        "PdfReader",
        lambda source: _FakePdfReader(
            [_FakePage("secret")],
            encrypted=True,
            decrypt_result=0,
        ),
    )

    with pytest.raises(EncryptedDocumentError):
        parse_pdf(b"%PDF-fake")


def test_pptx_parser_extracts_title_bullets_and_table() -> None:
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Qdrant"
    body = slide.placeholders[1].text_frame
    body.clear()
    body.paragraphs[0].text = "Dense vector"
    second = body.add_paragraph()
    second.text = "Sparse vector"
    second.level = 0
    table = slide.shapes.add_table(
        rows=2,
        cols=2,
        left=Inches(1),
        top=Inches(4),
        width=Inches(6),
        height=Inches(1),
    ).table
    table.cell(0, 0).text = "Key"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "DB"
    table.cell(1, 1).text = "Qdrant"
    output = BytesIO()
    presentation.save(output)

    deck = parse_pptx(output.getvalue(), deck_title="Backend")

    assert deck.source_format is SourceFormat.PPTX
    assert deck.title == "Backend"
    assert deck.slides[0].title == "Qdrant"
    assert any(
        block.kind is BlockKind.BULLET_GROUP
        and "Dense vector" in block.text
        and "Sparse vector" in block.text
        for block in deck.slides[0].blocks
    )
    table_blocks = [
        block for block in deck.slides[0].blocks if block.kind is BlockKind.TABLE_ROW_GROUP
    ]
    assert [block.text for block in table_blocks] == ["Key | Value", "DB | Qdrant"]
